"""Per-format text extraction handlers for the RAG ingestion pipeline.

Each function returns a ParsedDocument using the same types as ingestion.py.
"""
from __future__ import annotations

import csv
import io
import logging
import os
import re
from concurrent.futures import ThreadPoolExecutor, TimeoutError
from uuid import uuid4

from .ingestion import ParsedDocument, ParsedPage

logger = logging.getLogger("tablo-rag.parsers")
_VISION_TIMEOUT_S = float(os.getenv("RAG_VISION_TIMEOUT_S", "8"))


def _call_with_timeout(fn, timeout_s: float):
    if timeout_s <= 0:
        return fn()
    with ThreadPoolExecutor(max_workers=1) as executor:
        future = executor.submit(fn)
        return future.result(timeout=timeout_s)


def parse_docx(file_path: str) -> ParsedDocument:
    """Extract text from DOCX using python-docx."""
    from docx import Document

    doc_name = os.path.basename(file_path)
    doc = Document(file_path)
    pages: list[ParsedPage] = []
    offset = 0
    section_title = None

    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name and "Heading" in para.style.name:
            section_title = text
        start = offset
        end = offset + len(text)
        pages.append(ParsedPage(
            page_number=None,
            section_title=section_title,
            text=text,
            char_offset_start=start,
            char_offset_end=end,
        ))
        offset = end + 1

    return ParsedDocument(
        doc_id=str(uuid4()), doc_name=doc_name, format="docx",
        pages=pages, total_chars=offset,
    )


def parse_pptx(file_path: str) -> ParsedDocument:
    """Extract text from PPTX — one page per slide."""
    from pptx import Presentation

    doc_name = os.path.basename(file_path)
    prs = Presentation(file_path)
    pages: list[ParsedPage] = []
    offset = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        title = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_text_frame and shape.shape_id == slide.shapes.title_shape_id if hasattr(slide.shapes, 'title_shape_id') else False:
                title = shape.text_frame.text.strip()
        if not title and slide.shapes.title:
            title = slide.shapes.title.text.strip() if slide.shapes.title.has_text_frame else None

        page_text = " ".join(texts)
        if page_text.strip():
            start = offset
            end = offset + len(page_text)
            pages.append(ParsedPage(
                page_number=slide_num,
                section_title=title,
                text=page_text,
                char_offset_start=start,
                char_offset_end=end,
            ))
            offset = end + 1

    return ParsedDocument(
        doc_id=str(uuid4()), doc_name=doc_name, format="pptx",
        pages=pages, total_chars=offset,
    )


def parse_rtf(file_path: str) -> ParsedDocument:
    """Extract text from RTF using striprtf."""
    from striprtf.striprtf import rtf_to_text

    doc_name = os.path.basename(file_path)
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    text = rtf_to_text(raw).strip()

    paragraphs = re.split(r"\n{2,}", text)
    pages: list[ParsedPage] = []
    offset = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        start = offset
        end = offset + len(para)
        pages.append(ParsedPage(
            page_number=None, section_title=None,
            text=para, char_offset_start=start, char_offset_end=end,
        ))
        offset = end + 2

    return ParsedDocument(
        doc_id=str(uuid4()), doc_name=doc_name, format="rtf",
        pages=pages, total_chars=offset,
    )


def parse_image(file_path: str, genai_client) -> ParsedDocument:
    """Extract text/content from an image using Gemini vision."""
    if genai_client is None:
        doc_name = os.path.basename(file_path)
        text = f"[Image: {doc_name}]"
        pages = [ParsedPage(
            page_number=1, section_title=None,
            text=text, char_offset_start=0, char_offset_end=len(text),
        )]
        return ParsedDocument(
            doc_id=str(uuid4()), doc_name=doc_name, format=os.path.splitext(doc_name)[1].lstrip("."),
            pages=pages, total_chars=len(text),
        )

    from google.genai import types as genai_types

    doc_name = os.path.basename(file_path)
    with open(file_path, "rb") as f:
        image_bytes = f.read()

    ext = os.path.splitext(file_path)[1].lower().lstrip(".")
    mime_map = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg", "webp": "image/webp", "heif": "image/heif"}
    mime = mime_map.get(ext, "image/png")

    try:
        def _run():
            return genai_client.models.generate_content(
                model="gemini-2.5-flash",
                contents=[
                    genai_types.Part.from_bytes(data=image_bytes, mime_type=mime),
                    "Extract ALL text visible in this image. Also describe any diagrams, charts, or visual structures. Return the text content faithfully.",
                ],
            )

        response = _call_with_timeout(_run, _VISION_TIMEOUT_S)
        text = (response.text or "").strip()
    except TimeoutError:
        logger.warning("Image text extraction timed out after %.0fs", _VISION_TIMEOUT_S)
        text = f"[Image: {doc_name}]"
    except Exception as e:
        logger.warning("Image text extraction failed: %s", e)
        text = f"[Image: {doc_name}]"

    pages = [ParsedPage(
        page_number=1, section_title=None,
        text=text, char_offset_start=0, char_offset_end=len(text),
    )] if text else []

    return ParsedDocument(
        doc_id=str(uuid4()), doc_name=doc_name, format=ext,
        pages=pages, total_chars=len(text),
    )


def parse_xlsx(file_path: str) -> ParsedDocument:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook

    doc_name = os.path.basename(file_path)
    wb = load_workbook(file_path, read_only=True, data_only=True)
    pages: list[ParsedPage] = []
    offset = 0

    for sheet_num, sheet_name in enumerate(wb.sheetnames, 1):
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            row_text = " | ".join(c for c in cells if c)
            if row_text.strip():
                rows.append(row_text)
        page_text = "\n".join(rows)
        if page_text.strip():
            start = offset
            end = offset + len(page_text)
            pages.append(ParsedPage(
                page_number=sheet_num, section_title=sheet_name,
                text=page_text, char_offset_start=start, char_offset_end=end,
            ))
            offset = end + 1

    wb.close()
    return ParsedDocument(
        doc_id=str(uuid4()), doc_name=doc_name, format="xlsx",
        pages=pages, total_chars=offset,
    )


def parse_xls(file_path: str) -> ParsedDocument:
    """Extract text from XLS using xlrd."""
    import xlrd

    doc_name = os.path.basename(file_path)
    wb = xlrd.open_workbook(file_path)
    pages: list[ParsedPage] = []
    offset = 0

    for sheet_num in range(wb.nsheets):
        ws = wb.sheet_by_index(sheet_num)
        rows = []
        for row_idx in range(ws.nrows):
            cells = [str(ws.cell_value(row_idx, col)) for col in range(ws.ncols)]
            row_text = " | ".join(c for c in cells if c.strip())
            if row_text.strip():
                rows.append(row_text)
        page_text = "\n".join(rows)
        if page_text.strip():
            start = offset
            end = offset + len(page_text)
            pages.append(ParsedPage(
                page_number=sheet_num + 1, section_title=ws.name,
                text=page_text, char_offset_start=start, char_offset_end=end,
            ))
            offset = end + 1

    return ParsedDocument(
        doc_id=str(uuid4()), doc_name=doc_name, format="xls",
        pages=pages, total_chars=offset,
    )


def parse_csv_file(file_path: str, delimiter: str = ",") -> ParsedDocument:
    """Extract text from CSV/TSV."""
    doc_name = os.path.basename(file_path)
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f, delimiter=delimiter)
        rows = []
        for row in reader:
            row_text = " | ".join(c.strip() for c in row if c.strip())
            if row_text:
                rows.append(row_text)

    page_text = "\n".join(rows)
    pages = [ParsedPage(
        page_number=1, section_title=None,
        text=page_text, char_offset_start=0, char_offset_end=len(page_text),
    )] if page_text.strip() else []

    return ParsedDocument(
        doc_id=str(uuid4()), doc_name=doc_name,
        format="csv" if delimiter == "," else "tsv",
        pages=pages, total_chars=len(page_text),
    )


def parse_html(file_path: str) -> ParsedDocument:
    """Extract text from HTML using BeautifulSoup."""
    from bs4 import BeautifulSoup

    doc_name = os.path.basename(file_path)
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    # Remove script/style
    for tag in soup(["script", "style"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    title = soup.title.string.strip() if soup.title and soup.title.string else None

    pages = [ParsedPage(
        page_number=1, section_title=title,
        text=text, char_offset_start=0, char_offset_end=len(text),
    )] if text.strip() else []

    return ParsedDocument(
        doc_id=str(uuid4()), doc_name=doc_name, format="html",
        pages=pages, total_chars=len(text),
    )


def parse_doc(file_path: str, genai_client) -> ParsedDocument:
    """Extract text from DOC (old Word format). Try python-docx first, fallback to Gemini vision."""
    try:
        return parse_docx(file_path)
    except Exception:
        logger.info("python-docx failed for .doc file, falling back to Gemini vision")
        return parse_image(file_path, genai_client)


def parse_hwp(file_path: str, genai_client) -> ParsedDocument:
    """Extract text from HWP. Try hwp5 first, fallback to Gemini vision."""
    try:
        import hwp5
        # hwp5 support is limited — try basic text extraction
        from hwp5.hwp5html import open as hwp_open
        # This is unreliable, so we catch broadly
        raise ImportError("hwp5 text extraction not reliable")
    except Exception:
        logger.info("hwp5 not available or failed, falling back to Gemini vision for HWP")
        return parse_image(file_path, genai_client)
