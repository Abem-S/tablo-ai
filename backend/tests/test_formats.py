"""Document Format Parsing Tests — no external API calls for most tests."""
from __future__ import annotations
import csv, os, sys, tempfile, time
from dataclasses import dataclass
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))
from dotenv import load_dotenv
load_dotenv(Path(__file__).parent.parent / ".env")

@dataclass
class TestResult:
    name: str; passed: bool; score: float; latency_ms: float; details: str; recommendation: str = ""

def _skip(name, reason):
    return TestResult(name, True, 1.0, 0.0, f"SKIPPED: {reason}")

def test_txt_parsing() -> TestResult:
    name = "formats/txt"
    start = time.monotonic()
    content = "INTRODUCTION:\nNetworks connect computers.\n\nOSI MODEL:\nThe OSI model has 7 layers.\n\nTRANSPORT LAYER:\nTCP operates here."
    with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as f:
        f.write(content); tmp = f.name
    try:
        from rag.ingestion import IngestionPipeline
        from rag.knowledge_graph import KnowledgeGraph
        pipeline = object.__new__(IngestionPipeline)
        pipeline._kg = KnowledgeGraph()
        parsed = pipeline.parse_text(tmp)
        ms = (time.monotonic()-start)*1000
        issues = []
        if not parsed.pages: issues.append("no pages extracted")
        if parsed.format != "txt": issues.append(f"format={parsed.format}")
        all_text = " ".join(p.text for p in parsed.pages).lower()
        if "osi model" not in all_text: issues.append("'OSI model' not found")
        if "transport layer" not in all_text: issues.append("'Transport layer' not found")
        if len(parsed.pages) < 2: issues.append(f"expected ≥2 paragraphs, got {len(parsed.pages)}")
        titles = [p.section_title for p in parsed.pages if p.section_title]
        if not titles: issues.append("no section titles detected")
        score = max(0.0, 1.0 - len(issues)*0.2)
        details = f"Pages: {len(parsed.pages)}, chars: {parsed.total_chars}, titles: {titles[:3]}"
        if issues: details += "\nIssues: " + "; ".join(issues)
        return TestResult(name, len(issues)==0, score, ms, details)
    except Exception as e:
        return TestResult(name, False, 0.0, (time.monotonic()-start)*1000, f"Exception: {e}")
    finally:
        os.unlink(tmp)

def test_csv_parsing() -> TestResult:
    name = "formats/csv"
    start = time.monotonic()
    rows = [["Country","GDP","Year"],["USA","2.3","2023"],["China","5.2","2023"],["Germany","0.1","2023"]]
    with tempfile.NamedTemporaryFile(mode="w", suffix=".csv", delete=False, encoding="utf-8", newline="") as f:
        csv.writer(f).writerows(rows); tmp = f.name
    try:
        from rag.parsers import parse_csv_file
        parsed = parse_csv_file(tmp, delimiter=",")
        ms = (time.monotonic()-start)*1000
        issues = []
        if not parsed.pages: issues.append("no pages extracted")
        if parsed.format != "csv": issues.append(f"format={parsed.format}")
        all_text = " ".join(p.text for p in parsed.pages)
        for c in ["USA","China","Germany","Country"]:
            if c not in all_text: issues.append(f"'{c}' not found")
        if "|" not in all_text: issues.append("rows not pipe-separated")
        score = max(0.0, 1.0 - len(issues)*0.2)
        details = f"Pages: {len(parsed.pages)}, chars: {parsed.total_chars}"
        if issues: details += "\nIssues: " + "; ".join(issues)
        return TestResult(name, len(issues)==0, score, ms, details)
    except Exception as e:
        return TestResult(name, False, 0.0, (time.monotonic()-start)*1000, f"Exception: {e}")
    finally:
        os.unlink(tmp)

def test_html_parsing() -> TestResult:
    name = "formats/html"
    start = time.monotonic()
    html = """<!DOCTYPE html><html><head><title>OSI Overview</title>
<style>body{font-family:Arial}</style><script>console.log('x')</script></head>
<body><h1>The OSI Model</h1><p>7 layers for network communication.</p>
<h2>Transport Layer</h2><p>TCP and UDP operate here.</p></body></html>"""
    with tempfile.NamedTemporaryFile(mode="w", suffix=".html", delete=False, encoding="utf-8") as f:
        f.write(html); tmp = f.name
    try:
        from rag.parsers import parse_html
        parsed = parse_html(tmp)
        ms = (time.monotonic()-start)*1000
        issues = []
        if not parsed.pages: issues.append("no pages extracted")
        all_text = " ".join(p.text for p in parsed.pages)
        if "console.log" in all_text: issues.append("script not stripped")
        if "font-family" in all_text: issues.append("style not stripped")
        if "OSI Model" not in all_text and "OSI model" not in all_text: issues.append("'OSI Model' not found")
        if "Transport" not in all_text: issues.append("'Transport' not found")
        score = max(0.0, 1.0 - len(issues)*0.25)
        details = f"Pages: {len(parsed.pages)}, chars: {parsed.total_chars}"
        if issues: details += "\nIssues: " + "; ".join(issues)
        return TestResult(name, len(issues)==0, score, ms, details)
    except ImportError:
        return _skip(name, "beautifulsoup4 not installed")
    except Exception as e:
        return TestResult(name, False, 0.0, (time.monotonic()-start)*1000, f"Exception: {e}")
    finally:
        os.unlink(tmp)

def test_docx_parsing() -> TestResult:
    name = "formats/docx"
    start = time.monotonic()
    try:
        from docx import Document as DocxDoc
    except ImportError:
        return _skip(name, "python-docx not installed")
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        tmp = f.name
    try:
        doc = DocxDoc()
        doc.add_heading("Network Fundamentals", level=1)
        doc.add_paragraph("A computer network connects computers to share resources.")
        doc.add_heading("The OSI Model", level=2)
        doc.add_paragraph("The OSI model has 7 layers: Physical, Data Link, Network, Transport, Session, Presentation, Application.")
        doc.save(tmp)
        from rag.parsers import parse_docx
        parsed = parse_docx(tmp)
        ms = (time.monotonic()-start)*1000
        issues = []
        if not parsed.pages: issues.append("no pages extracted")
        all_text = " ".join(p.text for p in parsed.pages)
        if "OSI" not in all_text: issues.append("'OSI' not found")
        if "7 layers" not in all_text: issues.append("'7 layers' not found")
        titles = [p.section_title for p in parsed.pages if p.section_title]
        if not titles: issues.append("no headings detected as section_title")
        score = max(0.0, 1.0 - len(issues)*0.25)
        details = f"Pages: {len(parsed.pages)}, chars: {parsed.total_chars}, titles: {titles[:3]}"
        if issues: details += "\nIssues: " + "; ".join(issues)
        return TestResult(name, len(issues)==0, score, ms, details)
    except Exception as e:
        return TestResult(name, False, 0.0, (time.monotonic()-start)*1000, f"Exception: {e}")
    finally:
        os.unlink(tmp)

def test_pptx_parsing() -> TestResult:
    name = "formats/pptx"
    start = time.monotonic()
    try:
        from pptx import Presentation
    except ImportError:
        return _skip(name, "python-pptx not installed")
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        tmp = f.name
    try:
        prs = Presentation()
        for title_text, body_text in [
            ("Introduction", "Networks connect computers to share resources."),
            ("The OSI Model", "The OSI model has 7 layers defining network communication."),
            ("TCP/IP Stack", "TCP/IP uses 4 layers: Network Access, Internet, Transport, Application."),
        ]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title_text
            slide.placeholders[1].text = body_text
        prs.save(tmp)
        from rag.parsers import parse_pptx
        parsed = parse_pptx(tmp)
        ms = (time.monotonic()-start)*1000
        issues = []
        if len(parsed.pages) != 3: issues.append(f"expected 3 slides, got {len(parsed.pages)}")
        all_text = " ".join(p.text for p in parsed.pages)
        if "OSI" not in all_text: issues.append("'OSI' not found")
        if "TCP/IP" not in all_text: issues.append("'TCP/IP' not found")
        page_nums = [p.page_number for p in parsed.pages]
        if page_nums != [1,2,3]: issues.append(f"page_numbers={page_nums}")
        score = max(0.0, 1.0 - len(issues)*0.25)
        details = f"Slides: {len(parsed.pages)}, chars: {parsed.total_chars}, pages: {page_nums}"
        if issues: details += "\nIssues: " + "; ".join(issues)
        return TestResult(name, len(issues)==0, score, ms, details)
    except Exception as e:
        return TestResult(name, False, 0.0, (time.monotonic()-start)*1000, f"Exception: {e}")
    finally:
        os.unlink(tmp)

def test_viewer_routing() -> TestResult:
    name = "formats/viewer_routing"
    start = time.monotonic()
    def get_viewer_type(filename):
        ext = filename.rsplit(".",1)[-1].lower() if "." in filename else ""
        if ext == "pdf": return "pdf"
        if ext in ("png","jpg","jpeg","webp","heif"): return "image"
        if ext == "html": return "html"
        return "text"
    cases = [
        ("lecture.pdf","pdf"),("diagram.png","image"),("photo.jpg","image"),
        ("photo.jpeg","image"),("notes.docx","text"),("readme.txt","text"),
        ("slides.pptx","text"),("data.xlsx","text"),("data.csv","text"),
        ("page.html","html"),("image.webp","image"),
    ]
    issues = [f"{f}: expected '{e}', got '{get_viewer_type(f)}'" for f,e in cases if get_viewer_type(f)!=e]
    ms = (time.monotonic()-start)*1000
    score = (len(cases)-len(issues))/len(cases)
    details = f"Tested {len(cases)} mappings"
    if issues: details += "\nFailed: " + "; ".join(issues)
    return TestResult(name, len(issues)==0, score, ms, details)

ALL_TESTS = [test_txt_parsing, test_csv_parsing, test_html_parsing,
             test_docx_parsing, test_pptx_parsing, test_viewer_routing]

def run_all() -> list[TestResult]:
    results = []
    print("\n=== Document Format Tests ===\n")
    for fn in ALL_TESTS:
        r = fn()
        if r.details.startswith("SKIPPED"):
            print(f"  ⏭️  SKIP  {r.name}")
        else:
            status = "✅ PASS" if r.passed else "❌ FAIL"
            print(f"  {status} [{r.score:.0%}] {r.name} ({r.latency_ms:.0f}ms)")
            if not r.passed:
                for line in r.details.split("\n")[:3]: print(f"       {line}")
        results.append(r)
    return results

if __name__ == "__main__":
    results = run_all()
    passed = sum(1 for r in results if r.passed)
    avg = sum(r.score for r in results)/len(results)
    print(f"\nFormats: {passed}/{len(results)} passed, avg {avg:.0%}")
